from pptx import Presentation
import sys
from pathlib import Path

def convert_pptx_to_txt(pptx_path_str):
    """
    Extracts all text from a .pptx file and saves it to a .txt file
    in the 'pptx_export' directory.
    """
    pptx_path = Path(pptx_path_str)
    if not pptx_path.exists():
        print(f"Error: File not found at {pptx_path}", file=sys.stderr)
        sys.exit(1)

    output_dir = Path('pptx_export') # Changed to pptx_export
    output_dir.mkdir(exist_ok=True)

    try:
        presentation = Presentation(pptx_path)
        full_text = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = "".join(run.text for run in paragraph.runs)
                        if paragraph_text:
                            full_text.append(paragraph_text)
        
        txt_content = '\n'.join(full_text)
        
        txt_file_name = f"{pptx_path.stem}.txt"
        txt_path = output_dir / txt_file_name

        print(f"  -> Extracting text from '{pptx_path.name}' to '{txt_path}'")

        with open(txt_path, 'w', encoding='utf-8') as txtfile:
            txtfile.write(txt_content)

    except Exception as e:
        print(f"An error occurred: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python convert_pptx.py <path_to_pptx_file>", file=sys.stderr)
        sys.exit(1)
    
    file_to_convert = sys.argv[1]
    print(f"Extracting text from '{file_to_convert}'...")
    convert_pptx_to_txt(file_to_convert)
    print("Extraction complete.")
